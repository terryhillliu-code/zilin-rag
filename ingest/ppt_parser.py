"""
PPT 解析器 - 从 PowerPoint 文件提取文字和图片

支持：
- .pptx 格式（PowerPoint 2007+）
- 提取标题、正文、备注
- 提取嵌入图片
- 转换为 Markdown 格式
"""
import os
import re
import uuid
from typing import List, Optional, Dict, Any
from dataclasses import dataclass, field
from pathlib import Path
import logging

logger = logging.getLogger(__name__)


@dataclass
class PPTSlide:
    """PPT 幻灯片内容"""
    slide_num: int                          # 幻灯片编号（从 1 开始）
    title: str = ""                         # 标题
    content: List[str] = field(default_factory=list)  # 文本段落
    images: List[str] = field(default_factory=list)   # 图片路径
    notes: str = ""                         # 备注文字
    layout: str = ""                        # 布局类型


@dataclass
class PPTDocument:
    """PPT 文档"""
    file_path: str                          # 原始文件路径
    filename: str                           # 文件名
    slides: List[PPTSlide] = field(default_factory=list)
    total_slides: int = 0
    total_images: int = 0
    metadata: Dict[str, Any] = field(default_factory=dict)


class PPTParser:
    """
    PPT 解析器

    使用 python-pptx 提取文字和图片
    """

    def __init__(
        self,
        extract_images: bool = True,
        image_output_dir: Optional[str] = None,
        min_image_size: int = 100
    ):
        """
        初始化 PPT 解析器

        Args:
            extract_images: 是否提取图片
            image_output_dir: 图片输出目录
            min_image_size: 最小图片尺寸
        """
        self.extract_images = extract_images
        self.image_output_dir = Path(image_output_dir or "/tmp/zhiwei_ppt_images")
        self.min_image_size = min_image_size

        if extract_images:
            self.image_output_dir.mkdir(parents=True, exist_ok=True)

    def parse(self, file_path: str) -> PPTDocument:
        """
        解析 PPT 文件

        Args:
            file_path: PPT 文件路径

        Returns:
            PPTDocument 包含所有幻灯片内容
        """
        try:
            from pptx import Presentation
            from pptx.enum.shapes import MSO_SHAPE_TYPE
        except ImportError:
            raise ImportError("请安装 python-pptx: pip install python-pptx")

        if not os.path.exists(file_path):
            raise FileNotFoundError(f"文件不存在: {file_path}")

        ext = Path(file_path).suffix.lower()
        if ext not in [".pptx"]:
            raise ValueError(f"不支持的文件格式: {ext}，仅支持 .pptx")

        prs = Presentation(file_path)
        doc = PPTDocument(
            file_path=file_path,
            filename=Path(file_path).name,
            total_slides=len(prs.slides)
        )

        # 提取元数据
        doc.metadata = self._extract_metadata(prs)

        # 解析每张幻灯片
        total_images = 0
        for slide_num, slide in enumerate(prs.slides, 1):
            ppt_slide = self._parse_slide(slide, slide_num, Path(file_path).stem)
            doc.slides.append(ppt_slide)
            total_images += len(ppt_slide.images)

        doc.total_images = total_images
        logger.info(f"[PPTParser] 解析完成: {doc.total_slides} 张幻灯片, {total_images} 张图片")

        return doc

    def _extract_metadata(self, prs) -> Dict[str, Any]:
        """提取 PPT 元数据"""
        metadata = {}

        try:
            core_props = prs.core_properties
            metadata = {
                "author": core_props.author or "",
                "title": core_props.title or "",
                "subject": core_props.subject or "",
                "keywords": core_props.keywords or "",
                "created": str(core_props.created) if core_props.created else "",
                "modified": str(core_props.modified) if core_props.modified else "",
            }
        except Exception as e:
            logger.warning(f"[PPTParser] 提取元数据失败: {e}")

        return metadata

    def _parse_slide(self, slide, slide_num: int, filename: str) -> PPTSlide:
        """解析单张幻灯片"""
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        ppt_slide = PPTSlide(slide_num=slide_num)

        # 提取文本
        for shape in slide.shapes:
            # 获取标题
            if shape.has_text_frame:
                text = self._extract_text_from_shape(shape)

                # 判断是否为标题（通常是第一个文本框且字号较大）
                if not ppt_slide.title and shape == slide.shapes[0]:
                    ppt_slide.title = text.strip()
                elif text.strip():
                    ppt_slide.content.append(text.strip())

            # 提取图片
            if self.extract_images and shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image_path = self._save_image(shape, filename, slide_num)
                if image_path:
                    ppt_slide.images.append(image_path)

        # 提取备注
        if slide.has_notes_slide:
            notes_slide = slide.notes_slide
            notes_text = notes_slide.notes_text_frame.text
            if notes_text:
                ppt_slide.notes = notes_text.strip()

        return ppt_slide

    def _extract_text_from_shape(self, shape) -> str:
        """从形状中提取文本"""
        if not shape.has_text_frame:
            return ""

        texts = []
        for paragraph in shape.text_frame.paragraphs:
            para_text = []
            for run in paragraph.runs:
                para_text.append(run.text)
            if para_text:
                texts.append("".join(para_text))

        return "\n".join(texts)

    def _save_image(self, shape, filename: str, slide_num: int) -> Optional[str]:
        """保存图片到文件"""
        try:
            import hashlib

            image = shape.image
            image_bytes = image.blob

            # 计算哈希
            image_hash = hashlib.md5(image_bytes).hexdigest()[:8]

            # 生成文件名
            ext = image.ext
            output_name = f"{filename}_s{slide_num}_{image_hash}.{ext}"
            output_path = self.image_output_dir / output_name

            # 检查尺寸
            width_px = int(shape.width / 914400 * 96)
            height_px = int(shape.height / 914400 * 96)

            if width_px < self.min_image_size or height_px < self.min_image_size:
                return None

            # 保存
            with open(output_path, "wb") as f:
                f.write(image_bytes)

            return str(output_path)

        except Exception as e:
            logger.warning(f"[PPTParser] 保存图片失败: {e}")
            return None

    def to_markdown(self, doc: PPTDocument, include_images: bool = True) -> str:
        """
        将 PPT 转换为 Markdown

        Args:
            doc: PPT 文档
            include_images: 是否包含图片链接

        Returns:
            Markdown 格式文本
        """
        lines = []

        # 文档标题
        title = doc.metadata.get("title") or doc.filename
        lines.append(f"# {title}")
        lines.append("")

        # 元数据
        if doc.metadata.get("author"):
            lines.append(f"> 作者: {doc.metadata['author']}")
        if doc.metadata.get("keywords"):
            lines.append(f"> 关键词: {doc.metadata['keywords']}")
        lines.append("")

        # 各幻灯片
        for slide in doc.slides:
            # 幻灯片标题
            slide_title = slide.title or f"幻灯片 {slide.slide_num}"
            lines.append(f"## {slide.slide_num}. {slide_title}")
            lines.append("")

            # 正文内容
            for para in slide.content:
                if para != slide.title:
                    lines.append(para)
                    lines.append("")

            # 图片
            if include_images and slide.images:
                lines.append("**图片:**")
                for img_path in slide.images:
                    lines.append(f"![图片]({img_path})")
                lines.append("")

            # 备注
            if slide.notes:
                lines.append(f"> 备注: {slide.notes}")
                lines.append("")

            lines.append("---")
            lines.append("")

        return "\n".join(lines)

    def to_text(self, doc: PPTDocument) -> str:
        """
        将 PPT 转换为纯文本

        Args:
            doc: PPT 文档

        Returns:
            纯文本
        """
        lines = []

        for slide in doc.slides:
            # 标题
            if slide.title:
                lines.append(f"【{slide.slide_num}. {slide.title}】")

            # 内容
            for para in slide.content:
                if para != slide.title:
                    lines.append(para)

            # 备注
            if slide.notes:
                lines.append(f"[备注: {slide.notes}]")

            lines.append("")

        return "\n".join(lines)


def parse_ppt(file_path: str, **kwargs) -> PPTDocument:
    """
    便捷函数：解析 PPT 文件

    Args:
        file_path: PPT 文件路径
        **kwargs: 传递给 PPTParser 的参数

    Returns:
        PPTDocument
    """
    parser = PPTParser(**kwargs)
    return parser.parse(file_path)


def ppt_to_markdown(file_path: str, **kwargs) -> str:
    """
    便捷函数：PPT 转 Markdown

    Args:
        file_path: PPT 文件路径
        **kwargs: 传递给 PPTParser 的参数

    Returns:
        Markdown 文本
    """
    doc = parse_ppt(file_path, **kwargs)
    parser = PPTParser(**kwargs)
    return parser.to_markdown(doc)


if __name__ == "__main__":
    import argparse

    parser = argparse.ArgumentParser(description="PPT 解析测试")
    parser.add_argument("file", help="PPT 文件路径")
    parser.add_argument("--output", help="输出 Markdown 文件路径")
    parser.add_argument("--no-images", action="store_true", help="不提取图片")
    args = parser.parse_args()

    doc = parse_ppt(args.file, extract_images=not args.no_images)

    print(f"解析结果:")
    print(f"  文件: {doc.filename}")
    print(f"  幻灯片数: {doc.total_slides}")
    print(f"  图片数: {doc.total_images}")

    if args.output:
        ppt_parser = PPTParser()
        md = ppt_parser.to_markdown(doc)
        with open(args.output, "w", encoding="utf-8") as f:
            f.write(md)
        print(f"  已保存到: {args.output}")
    else:
        print("\n内容预览:")
        print(ppt_parser.to_text(doc)[:500])