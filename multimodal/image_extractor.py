"""
图片提取器 - 从 PDF/PPT 中提取图片

支持：
- PDF 图片提取（PyMuPDF）
- PPT 图片提取（python-pptx）
"""
import os
import hashlib
import uuid
from typing import List, Optional, Tuple
from dataclasses import dataclass
from pathlib import Path
import logging
import io

logger = logging.getLogger(__name__)


@dataclass
class ExtractedImage:
    """提取的图片"""
    path: str                    # 保存路径
    source_page: int             # 来源页码
    source_type: str             # 来源类型: pdf/ppt
    width: int = 0               # 宽度
    height: int = 0              # 高度
    format: str = ""             # 格式: png/jpg/...
    image_hash: str = ""         # 图片哈希（用于去重）


class ImageExtractor:
    """
    图片提取器

    从 PDF 和 PPT 文件中提取嵌入的图片
    """

    def __init__(self, output_dir: str = "/tmp/zhiwei_images"):
        """
        初始化图片提取器

        Args:
            output_dir: 图片输出目录
        """
        self.output_dir = Path(output_dir)
        self.output_dir.mkdir(parents=True, exist_ok=True)

    def extract_from_pdf(
        self,
        pdf_path: str,
        min_size: int = 100,
        dedupe: bool = True
    ) -> List[ExtractedImage]:
        """
        从 PDF 提取图片

        Args:
            pdf_path: PDF 文件路径
            min_size: 最小图片尺寸（像素），小于此值跳过
            dedupe: 是否去重

        Returns:
            提取的图片列表
        """
        try:
            import fitz  # PyMuPDF
        except ImportError:
            logger.error("[ImageExtractor] PyMuPDF 未安装")
            raise ImportError("请安装 PyMuPDF: pip install pymupdf")

        images = []
        seen_hashes = set()

        doc = fitz.open(pdf_path)
        filename = Path(pdf_path).stem

        for page_num in range(len(doc)):
            page = doc[page_num]
            image_list = page.get_images(full=True)

            for img_idx, img_info in enumerate(image_list):
                try:
                    # 获取图片
                    xref = img_info[0]
                    base_image = doc.extract_image(xref)

                    # 检查尺寸
                    width = base_image["width"]
                    height = base_image["height"]
                    if width < min_size or height < min_size:
                        continue

                    # 计算哈希（去重）
                    image_bytes = base_image["image"]
                    image_hash = hashlib.md5(image_bytes).hexdigest()[:16]

                    if dedupe and image_hash in seen_hashes:
                        continue
                    seen_hashes.add(image_hash)

                    # 保存图片
                    ext = base_image["ext"]
                    output_name = f"{filename}_p{page_num + 1}_{img_idx}_{image_hash}.{ext}"
                    output_path = self.output_dir / output_name

                    with open(output_path, "wb") as f:
                        f.write(image_bytes)

                    images.append(ExtractedImage(
                        path=str(output_path),
                        source_page=page_num + 1,
                        source_type="pdf",
                        width=width,
                        height=height,
                        format=ext,
                        image_hash=image_hash
                    ))

                except Exception as e:
                    logger.warning(f"[ImageExtractor] 提取图片失败 (页 {page_num}, 图 {img_idx}): {e}")
                    continue

        doc.close()
        logger.info(f"[ImageExtractor] 从 PDF 提取 {len(images)} 张图片: {pdf_path}")
        return images

    def extract_from_ppt(
        self,
        ppt_path: str,
        min_size: int = 100,
        dedupe: bool = True
    ) -> List[ExtractedImage]:
        """
        从 PPT 提取图片

        Args:
            ppt_path: PPT 文件路径
            min_size: 最小图片尺寸（像素）
            dedupe: 是否去重

        Returns:
            提取的图片列表
        """
        try:
            from pptx import Presentation
            from pptx.enum.shapes import MSO_SHAPE_TYPE
        except ImportError:
            logger.error("[ImageExtractor] python-pptx 未安装")
            raise ImportError("请安装 python-pptx: pip install python-pptx")

        images = []
        seen_hashes = set()

        prs = Presentation(ppt_path)
        filename = Path(ppt_path).stem

        for slide_num, slide in enumerate(prs.slides, 1):
            for shape_idx, shape in enumerate(slide.shapes):
                try:
                    # 检查是否为图片
                    if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
                        continue

                    # 获取图片数据
                    image = shape.image
                    image_bytes = image.blob

                    # 检查尺寸
                    width = shape.width
                    height = shape.height
                    # EMU 转像素 (1 inch = 914400 EMU, 1 inch = 96 px)
                    width_px = int(width / 914400 * 96)
                    height_px = int(height / 914400 * 96)

                    if width_px < min_size or height_px < min_size:
                        continue

                    # 计算哈希
                    image_hash = hashlib.md5(image_bytes).hexdigest()[:16]

                    if dedupe and image_hash in seen_hashes:
                        continue
                    seen_hashes.add(image_hash)

                    # 确定格式
                    ext = image.ext
                    output_name = f"{filename}_s{slide_num}_{shape_idx}_{image_hash}.{ext}"
                    output_path = self.output_dir / output_name

                    with open(output_path, "wb") as f:
                        f.write(image_bytes)

                    images.append(ExtractedImage(
                        path=str(output_path),
                        source_page=slide_num,
                        source_type="ppt",
                        width=width_px,
                        height=height_px,
                        format=ext,
                        image_hash=image_hash
                    ))

                except Exception as e:
                    logger.warning(f"[ImageExtractor] 提取图片失败 (幻灯片 {slide_num}, 形状 {shape_idx}): {e}")
                    continue

        logger.info(f"[ImageExtractor] 从 PPT 提取 {len(images)} 张图片: {ppt_path}")
        return images

    def extract(
        self,
        file_path: str,
        **kwargs
    ) -> List[ExtractedImage]:
        """
        自动检测文件类型并提取图片

        Args:
            file_path: 文件路径
            **kwargs: 传递给具体提取方法的参数

        Returns:
            提取的图片列表
        """
        ext = Path(file_path).suffix.lower()

        if ext == ".pdf":
            return self.extract_from_pdf(file_path, **kwargs)
        elif ext in [".ppt", ".pptx"]:
            return self.extract_from_ppt(file_path, **kwargs)
        else:
            logger.warning(f"[ImageExtractor] 不支持的文件类型: {ext}")
            return []

    def cleanup(self, max_age_hours: int = 24):
        """
        清理过期图片

        Args:
            max_age_hours: 最大保留时间（小时）
        """
        import time

        now = time.time()
        cleaned = 0

        for img_path in self.output_dir.glob("*"):
            if img_path.is_file():
                age_hours = (now - img_path.stat().st_mtime) / 3600
                if age_hours > max_age_hours:
                    img_path.unlink()
                    cleaned += 1

        if cleaned:
            logger.info(f"[ImageExtractor] 清理 {cleaned} 张过期图片")


if __name__ == "__main__":
    import argparse

    parser = argparse.ArgumentParser(description="图片提取测试")
    parser.add_argument("file", help="PDF 或 PPT 文件路径")
    parser.add_argument("--output", default="/tmp/zhiwei_images", help="输出目录")
    args = parser.parse_args()

    extractor = ImageExtractor(output_dir=args.output)
    images = extractor.extract(args.file)

    print(f"提取 {len(images)} 张图片:")
    for img in images:
        print(f"  - {img.path} ({img.width}x{img.height})")